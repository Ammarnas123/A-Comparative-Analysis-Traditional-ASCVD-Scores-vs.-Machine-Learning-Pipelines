from pptx import Presentation
from pptx.util import Inches

def create_detailed_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Predicting 10-Year Cardiovascular Disease Risk: A Machine Learning Approach"
    subtitle.text = "A Comparative Analysis: Traditional ASCVD Scores vs. Machine Learning Pipelines\nHealth Data Science Capstone Project\n[Your Name Here]"

    # --- Slide 2: Problem Statement & Clinical Context ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Problem Statement & Clinical Context"
    content.text_frame.text = "1. The Clinical Reality:"
    p = content.text_frame.add_paragraph()
    p.text = "Cardiovascular disease (CVD) is a leading global cause of mortality."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "2. The Limitation of Traditional Tools:"
    p = content.text_frame.add_paragraph()
    p.text = "Standardized calculators (like the Framingham Risk Score / ASCVD) use linear assumptions to predict a patient\u2019s 10-year risk."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "These generalized models often flag perfectly healthy patients as high-risk (false positives) or miss sub-populations with non-linear clinical interactions."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "3. Capstone Objective:"
    p = content.text_frame.add_paragraph()
    p.text = "Leverage 10-year longitudinal data to map the true deviation of traditional risk scores, and evaluate whether Machine Learning algorithms can achieve better clinical accuracy and recall."
    p.level = 1

    # --- Slide 3: Dataset Overview & EDA ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Dataset Overview & EDA"
    content.text_frame.text = "Data Source: The Framingham Heart Study Dataset."
    p = content.text_frame.add_paragraph()
    p.text = "Sample Size: 4,240 participants."
    p = content.text_frame.add_paragraph()
    p.text = "Target Variable: `TenYearCHD` (Binary outcome: 1 if the patient developed Coronary Heart Disease within 10 years)."
    p = content.text_frame.add_paragraph()
    p.text = "Exploratory Data Analysis (EDA) Insights:"
    p = content.text_frame.add_paragraph()
    p.text = "Class Imbalance: Only ~15% of the dataset (644 patients) actually developed CVD. The remaining ~85% remained healthy. This massive imbalance required specialized resampling techniques."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Features Covered: Demographics (Age, Sex), Behavioral (Smoking), Medical History (BP Meds, Diabetes), and Clinical Metrics (Total Cholesterol, Systolic BP, BMI, Glucose)."
    p.level = 1

    # --- Slide 4: Missing Value Imputation (KNN) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Data Preprocessing: KNN Imputation"
    content.text_frame.text = "Initial Challenge: The dataset contained 388 missing Glucose readings, 105 missing Education values, and 50 missing Total Cholesterol inputs."
    p = content.text_frame.add_paragraph()
    p.text = "Methodology: K-Nearest Neighbors (KNN) Imputation"
    p = content.text_frame.add_paragraph()
    p.text = "Instead of dropping valuable patient records or flatly inserting the dataset 'mean' (which can skew medical distributions), we used KNN."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Data was first scaled using `StandardScaler` to prevent variables with large values (like Glucose) from mathematically dominating distance calculations."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Patients with missing values were matched to their 5 most clinically similar neighbors (`n_neighbors=5`), predicting their exact glucose or cholesterol level based on physiological similarity."
    p.level = 1

    # --- Slide 5: The Traditional Clinical Baseline ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Traditional Clinical Baseline (ASCVD / FRS)"
    content.text_frame.text = "To establish a baseline, we calculated the 10-year continuous risk probabilities using an approximation of the Framingham Risk logistic coefficients."
    p = content.text_frame.add_paragraph()
    p.text = "Threshold: A >15% probability classified the patient as High Risk (CVD=1)."
    p = content.text_frame.add_paragraph()
    p.text = "Results: 'The Panicked Doctor Effect'"
    p = content.text_frame.add_paragraph()
    p.text = "Sensitivity (Recall): 99.84%. The traditional tool caught almost every single sick patient..."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Accuracy: 15.28%. ...but it did so by diagnosing 4,234 out of 4,240 perfectly healthy patients as High Risk (False Positives)."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Conclusion: While statistically 'safe' (no deaths missed), the 3,591 false alarms render it physically unusable in a modern hospital due to resource waste and patient trauma."
    p.level = 1

    # --- Slide 6: Class Imbalance & Train-Test Split ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Handling Class Imbalance & 70/30 Split"
    content.text_frame.text = "Data Splitting Strategy:"
    p = content.text_frame.add_paragraph()
    p.text = "A strictly stratified 70% Training / 30% Testing split was used."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Models are ONLY evaluated on the unseen 30% test set to measure real-world diagnostic performance."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "The Target Imbalance Problem:"
    p = content.text_frame.add_paragraph()
    p.text = "Because 85% of patients didn't get sick, standard algorithms try to artificially boost their own accuracy by blindly predicting 'Healthy' 100% of the time."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "The Solution: SMOTE"
    p = content.text_frame.add_paragraph()
    p.text = "Synthetic Minority Over-sampling Technique (SMOTE) was applied ONLY to the training set."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "This synthesizes artificial 'sick' patient data mathematically, forcing the models to learn what high-risk physiology looks like."
    p.level = 1

    # --- Slide 7: ML Baseline (KNN & Decision Trees) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "ML Baseline Testing: KNN & Decision Trees"
    content.text_frame.text = "Hyperparameter Tuning using GridSearchCV optimized for AUC-ROC (Distinguishing sick from healthy)."
    p = content.text_frame.add_paragraph()
    p.text = "1. K-Nearest Neighbors (KNN)"
    p = content.text_frame.add_paragraph()
    p.text = "GridSearchCV swept K=[3,5,7,9,11,15]."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Optimal K=5 achieved 64.9% Accuracy and 48.2% Recall. Drastically reduced false alarms but missed half the sick patients."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "2. Decision Tree Classifier"
    p = content.text_frame.add_paragraph()
    p.text = "Optimized for Cost-Complexity Pruning (Alpha). Tested [0.0, 0.001, 0.005, 0.01, 0.02]."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Optimal Alpha=0.005 prevented severe overfitting."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Result: 50.8% Accuracy, 75.6% Recall. The Decision Tree mathematically sacrificed overall accuracy to prioritize catching high-risk patients securely."
    p.level = 1

    # --- Slide 8: Advanced Ensemble Modeling ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Advanced Ensembles & Regression"
    content.text_frame.text = "1. Logistic Regression (The Clinical Standard)"
    p = content.text_frame.add_paragraph()
    p.text = "Weighted utilizing `class_weight='balanced'` instead of SMOTE for direct clinical coefficient interpretability."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Results: 66.8% Accuracy | 61.1% Recall. The most stable, middle-ground performer."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "2. Random Forest"
    p = content.text_frame.add_paragraph()
    p.text = "An ensemble of 100 trees restricted to depth 7. Resistant to overfitting."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Results: 74.1% Accuracy | 43.5% Recall."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "3. Gradient Boosting (GBM)"
    p = content.text_frame.add_paragraph()
    p.text = "A sequential ensemble that iteratively corrects previous tree errors."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Results: 79.6% Accuracy | 18.1% Recall."
    p.level = 1

    # --- Slide 9: Clinical Feature Engineering ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Clinical Feature Engineering"
    content.text_frame.text = "To un-block accuracy limits, we mathematically engineered 4 new physiological markers that are difficult for basic ML models to deduce on their own."
    p = content.text_frame.add_paragraph()
    p.text = "1. Pulse Pressure: (`sysBP - diaBP`). Indicates aortic stiffness."
    p = content.text_frame.add_paragraph()
    p.text = "2. Mean Arterial Pressure (MAP): Average blood pressure in a cardiac cycle."
    p = content.text_frame.add_paragraph()
    p.text = "3. Age x Cigarettes: Interaction term. Smoking 20 cigarettes a day at Age 65 is wildly different than at Age 30."
    p = content.text_frame.add_paragraph()
    p.text = "4. Age x sysBP: Interaction term mapping hypertension degradation over time."
    p = content.text_frame.add_paragraph()
    p.text = "Results of Engineering: When re-fed into the Decision Tree, accuracy jumped massively from 50.8% to 70.8%, but heavily degraded its 75% recall."

    # --- Slide 10: Neural Networks (ANN) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Artificial Neural Networks (MLP Classifier)"
    content.text_frame.text = "Deep Learning Architecture:"
    p = content.text_frame.add_paragraph()
    p.text = "A Multi-Layer Perceptron (MLP) was constructed with a hidden layer of 50 neurons."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Activation: 'ReLu' | Solver: 'Adam' (Adaptive Moment Estimation) | Regularization: L2 Alpha=0.01."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "The complete dataset (including the 4 new engineered features) was strictly scaled via `StandardScaler` to ensure swift gradient descent convergence."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Results: 73.35% Accuracy | 41.97% Recall."
    p = content.text_frame.add_paragraph()
    p.text = "Conclusion: Like the Random Forest, the ANN struggled with the small, imbalanced tabular data scale, preferring to increase accuracy by sacrificing minority class detection."
    p.level = 1

    # --- Slide 11: Ensemble Bagging Impact ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Ensemble Bagging (Bootstrap Aggregating)"
    content.text_frame.text = "Mechanism:"
    p = content.text_frame.add_paragraph()
    p.text = "Trained 100 independent Decision Trees on varied random subsets of the SMOTE-balanced training data (with replacement), taking a 'majority vote' for the final prediction."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Clinical Impact: High Accuracy vs. Dangerous Consensus"
    p = content.text_frame.add_paragraph()
    p.text = "Result: 81.05% Accuracy (The highest of all models tested) | 16.06% Recall (The lowest of all models tested)."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Compared to the single optimized Decision Tree, Accuracy jumped by +30.19%, but Recall crashed by -59.59%."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Interpretation: Bagging violently reduced variance/overfitting, achieving near-perfect specificity (no false alarms). However, the 'safe' majority vote of the 100 trees caused it to miss 84% of actual heart disease cases, proving single pruned trees are better for minority patient capture."
    p.level = 1

    # --- Slide 12: Feature Importance ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Feature Importance: What drives CVD?"
    content.text_frame.text = "The simple, highly pruned Decision tree (Alpha=0.005) completely dropped all complex features, proving that 10-year CVD risk can be overwhelmingly predicted by just two integers:"
    p = content.text_frame.add_paragraph()
    p.text = "1. Age (80.03% Importance in predicting CVD)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "2. Systolic Blood Pressure (19.97% Importance)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "The Random Forest ensemble agreed, strictly ranking the Top 5 Drivers of 10-Year Heart Disease as:"
    p = content.text_frame.add_paragraph()
    p.text = "1. Age (24.1%)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "2. Systolic Blood Pressure (14.6%)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "3. Total Cholesterol (9.5%)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "4. Diastolic Blood Pressure (9.4%)"
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "5. Fasting Glucose (8.5%)"
    p.level = 1

    # --- Slide 13: Comprehensive Summary Table ---
    slide_layout_title_only = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout_title_only)
    title = slide.shapes.title
    title.text = "Comprehensive Model Performance Table"
    from pptx.util import Inches
    table_data = [
        ["Technique", "Accuracy", "Sensitivity (Recall)", "Specificity"],
        ["Traditional ASCVD", "15.28%", "99.84%", "0.14%"],
        ["Decision Tree (Tuned)", "50.86%", "75.65%", "46.43%"],
        ["Engineered DT", "70.83%", "57.51%", "73.20%"],
        ["Logistic Regression", "66.82%", "61.14%", "67.84%"],
        ["Random Forest", "74.14%", "43.52%", "79.61%"],
        ["Gradient Boosting", "79.56%", "18.13%", "90.55%"],
        ["Bagging (100 Trees)", "81.05%", "16.06%", "92.65%"],
        ["Neural Network (ANN)","73.35%", "41.97%", "78.96%"]
    ]
    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table
    for col_idx, col_name in enumerate(table_data[0]):
        table.cell(0, col_idx).text = col_name
        table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True
    for row_idx in range(1, rows):
        for col_idx in range(cols):
            table.cell(row_idx, col_idx).text = table_data[row_idx][col_idx]

    # --- Slide 14: Clinical Trade-offs ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "The Clinical Trade-Off"
    content.text_frame.text = "The analysis revealed a strict inverse relationship between Model Accuracy and Patient Safety (Recall)."
    p = content.text_frame.add_paragraph()
    p.text = "The 'Dangerous Perfectionists': Bagging (81%) and Gradient Boosting (79.5%) achieved the highest overall test accuracy and massive Specificity (90%+), practically eliminating false alarms."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "However, they accomplished this by only catching ~16-18% of actual heart disease patients, acting far too conservatively for clinical deployment."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "The 'Safety Net': The perfectly tuned Decision Tree (Alpha 0.005) acted as the inverse. By heavily pruning feature noise, it ignored overall accuracy (50.8%) to successfully capture 75.6% of patients destined for heart disease."
    p.level = 1

    # --- Slide 15: Conclusion & Final Recommendations ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Final Conclusion & Recommendations"
    content.text_frame.text = "Conclusion: Machine Learning definitively solves the 'Panicked Doctor' problem of Traditional ASCVD estimators, raising diagnostic precision mathematically from 15% up to ~81%."
    p = content.text_frame.add_paragraph()
    p.text = "Final Model Selection depends heavily on the intended Hospital deployment strategy:"
    p = content.text_frame.add_paragraph()
    p.text = "Recommendation 1 (First-Line Triage/Screening Tool): Deploy the highly pruned Decision Tree. It acts as an aggressive safety net, prioritizing a 75% patient capture rate perfectly suited for preventative cardiology."
    p.level = 1
    p = content.text_frame.add_paragraph()
    p.text = "Recommendation 2 (Diagnostic Assistant): Employ Logistic Regression. With its stable ~67% Accuracy and ~61% Recall, it acts as a highly interpretable 'middle-ground' tool for doctors looking for a statistically balanced second opinion."
    p.level = 1

    # Save the presentation
    save_path = r"C:\Users\HP\Downloads\AML- CapstoneProject\Framingham_CVD assesment\Capstone_Final_Presentation_Detailed_v2.pptx"
    prs.save(save_path)
    print(f"✅ Detailed PowerPoint presentation successfully created at:\n{save_path}")

if __name__ == "__main__":
    create_detailed_presentation()